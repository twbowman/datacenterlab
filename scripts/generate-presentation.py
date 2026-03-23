#!/usr/bin/env python3
"""Generate PowerPoint presentation from the presentation guide."""

import subprocess
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# --- Color palette (dark professional theme) ---
BG_DARK = RGBColor(0x1B, 0x1B, 0x2F)
BG_SECTION = RGBColor(0x16, 0x21, 0x3E)
ACCENT = RGBColor(0x00, 0xD2, 0xFF)
ACCENT2 = RGBColor(0x7C, 0x3A, 0xED)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY = RGBColor(0x99, 0x99, 0x99)
TABLE_HEADER_BG = RGBColor(0x00, 0xD2, 0xFF)
TABLE_ROW_BG = RGBColor(0x22, 0x2B, 0x45)
TABLE_ALT_BG = RGBColor(0x1B, 0x23, 0x3D)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=18,
    color=WHITE,
    bold=False,
    alignment=PP_ALIGN.LEFT,
    font_name="Calibri",
):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide(
    slide,
    bullets,
    start_top=None,
    font_size=18,
    color=LIGHT_GRAY,
    left=None,
    width=None,
):
    if start_top is None:
        start_top = Inches(1.8)
    if left is None:
        left = Inches(0.8)
    if width is None:
        width = Inches(11.5)
    txBox = slide.shapes.add_textbox(left, start_top, width, Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)
        p.level = 0
    return txBox


def add_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def render_mermaid_png(mermaid_code, output_path):
    """Render mermaid code directly to PNG using mmdc CLI."""
    try:
        with tempfile.NamedTemporaryFile(mode="w", suffix=".mmd", delete=False) as f:
            f.write(mermaid_code)
            mmd_path = f.name
        for cmd in [
            [
                "mmdc",
                "-i",
                mmd_path,
                "-o",
                str(output_path),
                "-b",
                "transparent",
                "-t",
                "dark",
                "-w",
                "1600",
            ],
            [
                "npx",
                "--yes",
                "@mermaid-js/mermaid-cli",
                "-i",
                mmd_path,
                "-o",
                str(output_path),
                "-b",
                "transparent",
                "-t",
                "dark",
                "-w",
                "1600",
            ],
        ]:
            try:
                result = subprocess.run(cmd, capture_output=True, timeout=60)
                if result.returncode == 0 and Path(output_path).exists():
                    return True
            except (FileNotFoundError, subprocess.TimeoutExpired):
                continue
    except Exception:
        pass
    return False


def add_diagram_or_placeholder(slide, mermaid_code, label, left, top, width, height):
    """Try to render mermaid diagram as PNG image, fall back to placeholder box."""
    with tempfile.TemporaryDirectory() as tmpdir:
        png_path = Path(tmpdir) / "diagram.png"

        if render_mermaid_png(mermaid_code, png_path):
            slide.shapes.add_picture(str(png_path), left, top, width, height)
            return

    # Fallback: styled placeholder box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x22, 0x2B, 0x45)
    shape.line.color.rgb = ACCENT
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = f"[{label}]"
    p.font.size = Pt(14)
    p.font.color.rgb = ACCENT
    p.font.name = "Calibri"
    p2 = tf.add_paragraph()
    p2.text = "Render with Mermaid Live Editor or mmdc CLI"
    p2.font.size = Pt(11)
    p2.font.color.rgb = MID_GRAY
    p2.alignment = PP_ALIGN.CENTER


def add_table(slide, headers, rows, left, top, col_widths):
    """Add a styled table to the slide."""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_width = sum(col_widths)
    row_height = Inches(0.45)
    table_height = row_height * num_rows

    shape = slide.shapes.add_table(num_rows, num_cols, left, top, table_width, table_height)
    table = shape.table

    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = BG_DARK
            paragraph.font.name = "Calibri"
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG

    # Data rows
    for r, row in enumerate(rows):
        bg = TABLE_ROW_BG if r % 2 == 0 else TABLE_ALT_BG
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.color.rgb = LIGHT_GRAY
                paragraph.font.name = "Calibri"
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Use blank layout
    blank_layout = prs.slide_layouts[6]

    # =========================================================================
    # SLIDE 1: Title
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)
    add_text_box(
        slide,
        Inches(1),
        Inches(2),
        Inches(11),
        Inches(1.2),
        "Production Network Testing Lab",
        font_size=44,
        color=WHITE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_text_box(
        slide,
        Inches(1),
        Inches(3.4),
        Inches(11),
        Inches(0.8),
        "Multi-Vendor Datacenter Automation with gNMI",
        font_size=24,
        color=ACCENT,
        alignment=PP_ALIGN.CENTER,
    )
    add_text_box(
        slide,
        Inches(1),
        Inches(4.5),
        Inches(11),
        Inches(0.6),
        "Deploy  •  Configure  •  Monitor  •  Validate",
        font_size=18,
        color=LIGHT_GRAY,
        alignment=PP_ALIGN.CENTER,
    )
    add_notes(
        slide,
        "Elevator pitch: This is a production-grade datacenter network lab that lets you "
        "deploy, configure, monitor, and validate multi-vendor network fabrics using "
        "containerized switches. Everything built here works identically in production. "
        "The only thing that changes is the inventory file.",
    )

    # =========================================================================
    # SLIDE 2: The Problem
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)
    add_text_box(
        slide,
        Inches(0.8),
        Inches(0.5),
        Inches(11),
        Inches(0.8),
        "The Problem",
        font_size=36,
        color=ACCENT,
        bold=True,
    )
    bullets = [
        "• Datacenter network automation is fragmented — every vendor has its own CLI, API, data models, and metric formats",
        "• Testing network changes in production is risky — no safe place to iterate",
        "• Monitoring across vendors requires vendor-specific dashboards and queries",
        "• Configuration validation is manual and error-prone",
        "• Existing platforms that integrate these functions are typically vendor-locked, proprietary, or expensive",
    ]
    add_bullet_slide(slide, bullets, start_top=Inches(1.6))
    add_text_box(
        slide,
        Inches(0.8),
        Inches(5.5),
        Inches(11),
        Inches(0.8),
        '"How do you test a BGP configuration change across 4 different vendor types\nbefore pushing to production?"',
        font_size=20,
        color=ACCENT,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_notes(
        slide,
        "Key question to pose to the audience. This sets up the entire presentation. "
        "Most teams either don't test cross-vendor changes, or they test on one vendor "
        "and hope it works on the others.",
    )

    # =========================================================================
    # SLIDE 3: Unmet Needs in the Market
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)
    add_text_box(
        slide,
        Inches(0.8),
        Inches(0.5),
        Inches(11),
        Inches(0.8),
        "Unmet Needs in the Market",
        font_size=36,
        color=ACCENT,
        bold=True,
    )
    bullets = [
        "• Multi-vendor gNMI automation — very little reference material for cross-vendor gNMI config with a single dispatcher pattern",
        "• Telemetry normalization — mapping vendor-specific metric names to universal names barely exists as reusable examples",
        "• Validation as code — industry embraced config-as-code, but validation is still 'SSH in and check'",
        "• Laptop-runnable multi-vendor lab — most Containerlab examples are single-vendor or topology-only",
    ]
    add_bullet_slide(slide, bullets, start_top=Inches(1.6), font_size=20)
    add_text_box(
        slide,
        Inches(0.8),
        Inches(5.2),
        Inches(11),
        Inches(1),
        "Not competing with CloudVision, Apstra, or NSO.\n"
        "Showing what's possible with open-source tools and gNMI. The integration across vendors is the value.",
        font_size=16,
        color=MID_GRAY,
        alignment=PP_ALIGN.CENTER,
    )
    add_notes(
        slide,
        "This project isn't a new tool — it's a working, integrated, open-source reference architecture. "
        "The individual tools (Ansible, Grafana, Prometheus, gNMIc) aren't novel — the integration across vendors is. "
        "Vendor docs show single-vendor examples. This project shows the cross-vendor reality including the ugly parts "
        "like SR Linux rate limits and OpenConfig gaps.",
    )

    # =========================================================================
    # SLIDE 4: What This Project Is
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)
    add_text_box(
        slide,
        Inches(0.8),
        Inches(0.5),
        Inches(11),
        Inches(0.8),
        "What This Project Is",
        font_size=36,
        color=ACCENT,
        bold=True,
    )
    bullets = [
        "• 2 spine switches (route reflectors) + 4 leaf switches + 4 clients",
        "• Full CLOS fabric: OSPF underlay, iBGP overlay, EVPN/VXLAN",
        "• 4 vendors: Nokia SR Linux, Arista cEOS, SONiC, Juniper cRPD",
        "• Runs on a laptop via Containerlab + OrbStack (macOS ARM)",
        "• Everything is infrastructure-as-code, version controlled",
    ]
    add_bullet_slide(slide, bullets, start_top=Inches(1.6), font_size=22)
    add_text_box(
        slide,
        Inches(0.8),
        Inches(5.5),
        Inches(11),
        Inches(0.6),
        "Demo: ./lab start  →  containers come up in ~2 minutes",
        font_size=18,
        color=ACCENT,
    )
    add_notes(
        slide,
        "A containerized datacenter network running real network operating systems. "
        "Demo opportunity: Show ./lab start and ./lab status.",
    )

    # =========================================================================
    # SLIDE 5: Architecture
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)
    add_text_box(
        slide,
        Inches(0.8),
        Inches(0.3),
        Inches(11),
        Inches(0.8),
        "Architecture",
        font_size=36,
        color=ACCENT,
        bold=True,
    )

    arch_mermaid = """graph TD
    subgraph "Automation"
        Ansible["Ansible + Dispatcher Pattern"]
    end
    subgraph "Network Fabric"
        SRL[SR Linux]
        EOS[Arista cEOS]
        SONiC[SONiC]
        JUN[Juniper cRPD]
    end
    subgraph "Telemetry Pipeline"
        gNMIc --> Normalize["Metric Normalization"] --> Prometheus --> Grafana
    end
    subgraph "Validation"
        gnmi_validate["gnmi_validate module"]
    end
    Ansible -->|gNMI SET| SRL & EOS & SONiC & JUN
    SRL & EOS & SONiC & JUN -->|gNMI SUBSCRIBE| gNMIc
    gnmi_validate -->|gNMI GET| SRL & EOS & SONiC & JUN"""

    add_diagram_or_placeholder(
        slide, arch_mermaid, "Architecture Diagram", Inches(1.5), Inches(1.3), Inches(10), Inches(5)
    )

    bullets_right = [
        "• gNMI: single transport for config, telemetry, validation",
        "• OpenConfig for telemetry, native YANG for config",
        "• Dispatcher pattern auto-routes to vendor-specific roles",
    ]
    add_bullet_slide(slide, bullets_right, start_top=Inches(6.2), font_size=14, color=MID_GRAY)
    add_notes(
        slide,
        "gNMI is the single transport for config (SET), telemetry (SUBSCRIBE), and validation (GET). "
        "OpenConfig for telemetry gives vendor-neutral dashboards. Native YANG for config gives full feature access. "
        "Ansible dispatcher pattern auto-routes to vendor-specific roles based on OS detection. "
        "Everything is infrastructure-as-code, version controlled.",
    )

    # =========================================================================
    # SLIDE 6: Multi-Vendor Dispatcher Pattern
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)
    add_text_box(
        slide,
        Inches(0.8),
        Inches(0.5),
        Inches(11),
        Inches(0.8),
        "Multi-Vendor Dispatcher Pattern",
        font_size=36,
        color=ACCENT,
        bold=True,
    )
    bullets = [
        "1. Dynamic inventory queries gNMI capabilities from each device",
        "2. OS detected automatically (srlinux, eos, sonic, junos)",
        "3. site.yml uses conditionals to route to vendor-specific roles",
        "4. Same playbook, same variables, different vendors",
    ]
    add_bullet_slide(slide, bullets, start_top=Inches(1.6), font_size=20)

    # Code snippet box
    code_text = (
        "- name: Configure SR Linux\n"
        "  include_role: { name: gnmi_interfaces }\n"
        "  when: ansible_network_os == 'nokia.srlinux'\n\n"
        "- name: Configure Arista\n"
        "  include_role: { name: eos_interfaces }\n"
        "  when: ansible_network_os == 'arista.eos'"
    )
    code_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(4), Inches(10), Inches(2.5)
    )
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(0x0D, 0x11, 0x17)
    code_box.line.color.rgb = ACCENT2
    code_box.line.width = Pt(1)
    tf = code_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code_text
    p.font.size = Pt(15)
    p.font.color.rgb = RGBColor(0xA0, 0xE8, 0xAF)
    p.font.name = "Courier New"

    add_notes(
        slide,
        "This is the core automation innovation. Walk through how it works. "
        "Demo opportunity: Run ansible-playbook -i inventory.yml site.yml and show it "
        "configuring interfaces, LLDP, OSPF, BGP, EVPN across all devices.",
    )

    # =========================================================================
    # SLIDE 7: EVPN/VXLAN Fabric
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)
    add_text_box(
        slide,
        Inches(0.8),
        Inches(0.5),
        Inches(11),
        Inches(0.8),
        "EVPN/VXLAN Fabric",
        font_size=36,
        color=ACCENT,
        bold=True,
    )
    bullets = [
        "• OSPF provides underlay reachability between loopbacks",
        "• iBGP with route reflectors (spines) distributes EVPN routes",
        "• VXLAN tunnels extend L2 domains across the fabric",
        "• 5 tenant VLANs mapped to VNIs (10010–10050), 2 L3 VRFs",
        "• Data model in group_vars/leafs.yml — change YAML, re-run playbook, fabric reconfigures",
    ]
    add_bullet_slide(slide, bullets, start_top=Inches(1.6), font_size=20)
    add_text_box(
        slide,
        Inches(0.8),
        Inches(5.5),
        Inches(11),
        Inches(0.6),
        "Demo: ping 10.10.100.12 — traffic traverses leaf1 → spine → leaf2 via VXLAN",
        font_size=18,
        color=ACCENT,
    )
    add_notes(
        slide,
        "Explain the overlay architecture. Demo opportunity: "
        "docker exec clab-gnmi-clos-client1 ping 10.10.100.12 — "
        "traffic traverses leaf1 → spine → leaf2 via VXLAN.",
    )

    # =========================================================================
    # SLIDE 8: Telemetry and Monitoring
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)
    add_text_box(
        slide,
        Inches(0.8),
        Inches(0.5),
        Inches(11),
        Inches(0.8),
        "Telemetry & Monitoring",
        font_size=36,
        color=ACCENT,
        bold=True,
    )

    # Left column: approach
    left_bullets = [
        "Two-tier telemetry:",
        "  • Tier 1: OpenConfig paths (vendor-neutral)",
        "  • Tier 2: Native vendor paths (OSPF, EVPN)",
        "",
        "Normalization pipeline:",
        "  • gNMIc collects via gNMI SUBSCRIBE",
        "  • Event processors → universal metric names",
        "  • network_interface_in_octets works for all vendors",
    ]
    add_bullet_slide(
        slide,
        left_bullets,
        start_top=Inches(1.6),
        font_size=18,
        left=Inches(0.8),
        width=Inches(5.5),
    )

    # Right column: dashboards
    right_bullets = [
        "9 Grafana dashboards:",
        "  • Universal Interfaces",
        "  • Universal BGP",
        "  • Universal LLDP",
        "  • Interface Performance",
        "  • BGP / OSPF / EVPN Stability",
        "  • Network Congestion",
        "  • Vendor SR Linux (native drill-down)",
    ]
    add_bullet_slide(
        slide,
        right_bullets,
        start_top=Inches(1.6),
        font_size=18,
        left=Inches(6.8),
        width=Inches(5.5),
    )

    add_text_box(
        slide,
        Inches(0.8),
        Inches(6.2),
        Inches(11),
        Inches(0.6),
        "Demo: Open Grafana at localhost:3000 — Universal Interfaces dashboard with all vendors",
        font_size=16,
        color=ACCENT,
    )
    add_notes(
        slide,
        "Two-tier telemetry approach. Tier 1 OpenConfig for vendor-neutral metrics. "
        "Tier 2 native vendor paths for features OpenConfig doesn't cover. "
        "Demo: Open Grafana, show Universal Interfaces dashboard, then drill into vendor-specific view.",
    )

    # =========================================================================
    # SLIDE 9: Validation Framework
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)
    add_text_box(
        slide,
        Inches(0.8),
        Inches(0.5),
        Inches(11),
        Inches(0.8),
        "Validation Framework",
        font_size=36,
        color=ACCENT,
        bold=True,
    )
    bullets = [
        "Custom gnmi_validate Ansible module:",
        "  • pygnmi queries device state via gNMI GET",
        "  • Compares actual state vs expected state from inventory",
        "  • Supports OpenConfig + vendor-native YANG (origin field)",
        "  • Structured pass/fail with diffs and remediation hints",
        "",
        "Validation playbooks:",
        "  • validate-bgp.yml — BGP sessions ESTABLISHED?",
        "  • validate-evpn.yml — EVPN routes advertised/received?",
        "  • validate-lldp.yml — LLDP neighbors match topology?",
        "  • validate-interfaces.yml — interfaces oper-up?",
    ]
    add_bullet_slide(slide, bullets, start_top=Inches(1.6), font_size=18)
    add_text_box(
        slide,
        Inches(0.8),
        Inches(5.8),
        Inches(11),
        Inches(0.8),
        "Demo: Run validation → break something → re-validate → see failure + remediation hint",
        font_size=16,
        color=ACCENT,
    )
    add_notes(
        slide,
        "Custom gnmi_validate Ansible module uses pygnmi to query device state via gNMI GET. "
        "Compares actual state against expected state from inventory variables. "
        "Demo: Run ansible-playbook playbooks/validate.yml, show structured output. "
        "Then intentionally break something (shut an interface) and re-run to show failure detection.",
    )

    return prs


def add_remaining_slides(prs):
    blank_layout = prs.slide_layouts[6]

    # =========================================================================
    # SLIDE 10: Production Portability
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)
    add_text_box(
        slide,
        Inches(0.8),
        Inches(0.5),
        Inches(11),
        Inches(0.8),
        "Production Portability",
        font_size=36,
        color=ACCENT,
        bold=True,
    )
    add_text_box(
        slide,
        Inches(0.8),
        Inches(1.4),
        Inches(11),
        Inches(0.6),
        "Everything built here transfers directly to production",
        font_size=22,
        color=WHITE,
    )

    headers = ["Component", "Lab", "Production", "What Changes"]
    rows = [
        ["Ansible playbooks", "Same", "Same", "Nothing"],
        ["Grafana dashboards", "Same", "Same", "Nothing"],
        ["gNMI subscriptions", "Same", "Same", "Nothing"],
        ["Validation checks", "Same", "Same", "Nothing"],
        ["Inventory file", "containerlab IPs", "datacenter IPs", "Only this"],
    ]
    col_widths = [Inches(3), Inches(2.5), Inches(2.5), Inches(2.5)]
    add_table(slide, headers, rows, Inches(1.2), Inches(2.2), col_widths)

    add_text_box(
        slide,
        Inches(0.8),
        Inches(5.8),
        Inches(11),
        Inches(0.6),
        "develop in lab  →  validate  →  update inventory  →  run same playbook against production",
        font_size=18,
        color=ACCENT,
        alignment=PP_ALIGN.CENTER,
    )
    add_notes(
        slide,
        "This is the key differentiator. Migration path: develop in lab, validate, "
        "update inventory (device IPs/credentials), run same Ansible playbook against production. "
        "Same dashboards, same queries, same validation checks.",
    )

    # =========================================================================
    # SLIDE 11: Engineering Challenges
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)
    add_text_box(
        slide,
        Inches(0.8),
        Inches(0.5),
        Inches(11),
        Inches(0.8),
        "Engineering Challenges",
        font_size=36,
        color=ACCENT,
        bold=True,
    )

    challenges = [
        (
            "SR Linux gNMI Rate Limit",
            "60 connections/minute per device. Solved by batching gNMI SET operations "
            "into single calls with multiple update-path/update-value pairs using Jinja2 loops. "
            "Reduced ~50 connections per host to ~15.",
        ),
        (
            "OpenConfig Gaps",
            "SR Linux doesn't expose OSPF or EVPN via OpenConfig. Solution: dual-schema approach "
            "using the gNMI origin field — OpenConfig for telemetry, native YANG for config.",
        ),
        (
            "Metric Normalization at Scale",
            "4 vendors × different path formats × different naming conventions. Built a normalization "
            "pipeline in gNMIc that produces identical metric names regardless of source vendor.",
        ),
    ]

    y = Inches(1.6)
    for title, desc in challenges:
        add_text_box(
            slide,
            Inches(0.8),
            y,
            Inches(11),
            Inches(0.4),
            title,
            font_size=20,
            color=ACCENT,
            bold=True,
        )
        add_text_box(
            slide,
            Inches(0.8),
            y + Inches(0.45),
            Inches(11),
            Inches(1),
            desc,
            font_size=16,
            color=LIGHT_GRAY,
        )
        y += Inches(1.7)

    add_notes(
        slide,
        "Interesting problems solved along the way. The SR Linux rate limit was the trickiest — "
        "60 connections/minute for CPU protection. We batch all operations into single gnmic calls. "
        "OpenConfig gaps required a dual-schema approach using the gNMI origin field.",
    )

    # =========================================================================
    # SLIDE 12: Building with AI
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)
    add_text_box(
        slide,
        Inches(0.8),
        Inches(0.5),
        Inches(11),
        Inches(0.8),
        "Building with AI",
        font_size=36,
        color=ACCENT,
        bold=True,
    )

    # Big number callout
    add_text_box(
        slide,
        Inches(0.8),
        Inches(1.5),
        Inches(5),
        Inches(0.7),
        "18 weeks estimated",
        font_size=32,
        color=MID_GRAY,
        bold=True,
    )
    add_text_box(
        slide,
        Inches(0.8),
        Inches(2.2),
        Inches(5),
        Inches(0.7),
        "~4 days actual",
        font_size=40,
        color=ACCENT,
        bold=True,
    )
    add_text_box(
        slide,
        Inches(0.8),
        Inches(3.0),
        Inches(5),
        Inches(0.5),
        "50 commits  •  12–15 March 2026",
        font_size=16,
        color=MID_GRAY,
    )

    # Right column: key points
    right_bullets = [
        "• Cross-domain knowledge on tap — Ansible, gNMI,\n  YANG, 4 vendor CLIs, Python, Grafana, Docker",
        "• Faster iteration on vendor-specific quirks",
        "• Documentation stays current — docs updated\n  with every code change automatically",
        "• Spec-driven development — requirements → design\n  → tasks before writing code",
        "• Honest feedback loop — AI challenges assumptions",
        "• Clean commit discipline — conventional commits,\n  pre-commit hooks, push after every change",
    ]
    add_bullet_slide(
        slide,
        right_bullets,
        start_top=Inches(1.5),
        font_size=16,
        color=LIGHT_GRAY,
        left=Inches(6.5),
        width=Inches(6),
    )

    add_text_box(
        slide,
        Inches(0.8),
        Inches(6.2),
        Inches(11.5),
        Inches(0.6),
        "AI didn't replace network engineering expertise — it amplified execution across domains",
        font_size=18,
        color=WHITE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_notes(
        slide,
        "The design spec estimated 18 weeks of human effort across 9 phases. "
        "With AI-assisted development (Kiro), the project went from first commit to production-ready "
        "in ~4 calendar days with 50 commits. AI bridges knowledge gaps across Ansible, gNMI, YANG models, "
        "4 vendor CLIs, Jinja2, Python, Prometheus, Grafana, Docker, and Containerlab. "
        "Every change is committed with descriptive conventional commit messages. "
        "The git log reads like a changelog, not a wall of WIP. "
        "AI handles commit message formatting, pre-commit hook compliance, and push cycle automatically. "
        "Key takeaway: AI amplified the ability to execute across multiple technology domains simultaneously.",
    )

    # =========================================================================
    # SLIDE 13: CI/CD and Testing
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)
    add_text_box(
        slide,
        Inches(0.8),
        Inches(0.5),
        Inches(11),
        Inches(0.8),
        "CI/CD & Testing",
        font_size=36,
        color=ACCENT,
        bold=True,
    )

    # Pipeline stages
    stages = [
        ("1. Lint", "Ruff, Mypy, yamllint,\nShellCheck, ansible-lint"),
        ("2. Security", "Bandit, Trivy, Checkov,\nGitleaks"),
        ("3. Tests", "Unit + Property-based\n(Hypothesis)"),
    ]
    x = Inches(0.8)
    for title, desc in stages:
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(3.5), Inches(1.8)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = BG_SECTION
        box.line.color.rgb = ACCENT
        box.line.width = Pt(1.5)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.color.rgb = ACCENT
        p.font.bold = True
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(15)
        p2.font.color.rgb = LIGHT_GRAY
        p2.font.name = "Calibri"
        p2.alignment = PP_ALIGN.CENTER
        x += Inches(4)

    test_bullets = [
        "• Unit tests: deployment, configuration, telemetry, validation, state management",
        "• Property-based tests: state management invariants, telemetry properties",
        "• Integration tests: end-to-end workflows, multi-vendor, monitoring stack",
        "• Pre-commit hooks enforce all lint + security checks locally before push",
    ]
    add_bullet_slide(slide, test_bullets, start_top=Inches(4.2), font_size=17)
    add_notes(
        slide,
        "3-stage GitHub Actions pipeline. Pre-commit hooks enforce all lint + security checks "
        "locally before push. Test structure covers unit, property-based, and integration tests.",
    )

    # =========================================================================
    # SLIDE 14: Key Numbers
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)
    add_text_box(
        slide,
        Inches(0.8),
        Inches(0.5),
        Inches(11),
        Inches(0.8),
        "Key Numbers",
        font_size=36,
        color=ACCENT,
        bold=True,
    )

    headers = ["Metric", "Value"]
    rows = [
        ["Vendors supported", "4 (SR Linux, Arista, SONiC, Juniper)"],
        ["Network devices", "6 (2 spines + 4 leafs)"],
        ["Client nodes", "4"],
        ["Ansible roles", "26"],
        ["Grafana dashboards", "9"],
        ["Telemetry metrics", "~1,560 OpenConfig + native"],
        ["Validation checks", "5 categories"],
        ["CI pipeline stages", "3 (lint, security, tests)"],
        ["Tenant VLANs / L3 VRFs", "5 VLANs (VNI 10010–10050) / 2 VRFs"],
        ["Design estimate vs actual", "18 weeks → 4 days with AI"],
    ]
    col_widths = [Inches(4), Inches(6.5)]
    add_table(slide, headers, rows, Inches(1.2), Inches(1.5), col_widths)
    add_notes(slide, "Key metrics for the project. Good reference slide for Q&A.")

    # =========================================================================
    # SLIDE 15: Live Demo Sequence
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)
    add_text_box(
        slide,
        Inches(0.8),
        Inches(0.5),
        Inches(11),
        Inches(0.8),
        "Live Demo Sequence",
        font_size=36,
        color=ACCENT,
        bold=True,
    )

    demo_steps = [
        "1.  Show topology-srlinux.yml",
        "2.  Deploy: ./lab start",
        "3.  Show running containers: ./lab status",
        "4.  Configure fabric: ansible-playbook site.yml",
        "5.  Verify OSPF + BGP neighbors",
        "6.  Test client connectivity via EVPN/VXLAN (ping)",
        "7.  Show Grafana dashboards (localhost:3000)",
        "8.  Run validation playbook",
        "9.  Break something → re-validate → see failure + remediation",
        "10. Fix it → re-validate → all green",
    ]
    add_bullet_slide(slide, demo_steps, start_top=Inches(1.5), font_size=20)
    add_notes(
        slide,
        "If doing a live demo, follow this sequence. "
        "All commands need orb -m clab prefix on macOS ARM. "
        "The break-and-fix cycle is the most impactful demo moment — "
        "shut an interface, show failure detection with remediation hint, then fix it.",
    )

    # =========================================================================
    # SLIDE 16: Q&A
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)
    add_text_box(
        slide,
        Inches(1),
        Inches(2.5),
        Inches(11),
        Inches(1),
        "Questions?",
        font_size=48,
        color=ACCENT,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_text_box(
        slide,
        Inches(1),
        Inches(4),
        Inches(11),
        Inches(0.6),
        "Production Network Testing Lab",
        font_size=22,
        color=LIGHT_GRAY,
        alignment=PP_ALIGN.CENTER,
    )
    add_notes(
        slide,
        "Anticipated questions:\n\n"
        "Q: Why not Terraform/Nornir/Napalm?\n"
        "A: Those focus on config management. This is a complete platform — deployment, config, "
        "telemetry, monitoring, and validation in one stack. Ansible chosen for network ops adoption.\n\n"
        "Q: Why gNMI instead of NETCONF/RESTCONF?\n"
        "A: gNMI provides streaming telemetry (not just config), uses efficient gRPC/protobuf, "
        "and is the industry direction. One protocol for SET, GET, and SUBSCRIBE.\n\n"
        "Q: Does this scale to production?\n"
        "A: Yes. Same playbooks, dashboards, subscriptions. Prometheus and gNMIc support horizontal scaling. "
        "Only the inventory file changes.\n\n"
        "Q: Why OpenConfig for telemetry but native YANG for config?\n"
        "A: OpenConfig telemetry gives vendor-neutral dashboards. But OpenConfig config support is incomplete — "
        "SR Linux doesn't support EVPN via OpenConfig. Native YANG gives full feature access.\n\n"
        "Q: What about the SR Linux rate limit?\n"
        "A: 60 gNMI connections/minute for CPU protection. We batch operations into single gnmic calls "
        "using Jinja2 loops. Reduced from ~50 to ~15 connections per host.",
    )

    return prs


def main():
    prs = build_presentation()
    add_remaining_slides(prs)

    output_path = Path("docs/presentation.pptx")
    prs.save(str(output_path))
    print(f"Presentation saved to {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
